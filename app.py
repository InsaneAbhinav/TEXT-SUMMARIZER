import streamlit as st
import spacy
from transformers import pipeline
from txtai.pipeline import Summary
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation  # Importing the library to handle PPT files
import pytesseract
from PIL import Image
import re
import base64
import difflib
import time  # Import the time module

# Set page layout
st.set_page_config(layout="wide")

# Specify the path to tesseract.exe for Windows (update this path for your system)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  # Windows

# Load summarization model
@st.cache_resource
def load_summary_model():
    return Summary()

# Text summarization
@st.cache_data
def text_summary(text, maxlength=None):
    summary = load_summary_model()
    result = summary(text, maxlength=maxlength)
    return result

# Extract text from PDF
def extract_text_from_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page_num in range(len(reader.pages)):
        text += reader.pages[page_num].extract_text()
    return text if text else "No readable text found in the PDF."

# Extract text from Word
def extract_text_from_word(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

# Extract text from PowerPoint slides
def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text if text else "No readable text found in the PPT."

# Extract text from image using OCR with error handling
def extract_text_from_image(file):
    try:
        image = Image.open(file)
        text = pytesseract.image_to_string(image)
        return text if text else "No readable text found in the image."
    except Exception as e:
        return f"Error in extracting text from image: {str(e)}"

# Preprocess text (optional, used to clean text)
def preprocess_text(text):
    text = re.sub(r'\s+', ' ', text)  # Remove extra spaces
    text = re.sub(r'[^A-Za-z0-9\s]', '', text)  # Remove special characters
    return text

# Download summary as .txt
def download_summary(summary_text, filename="summary.txt"):
    b64 = base64.b64encode(summary_text.encode()).decode()
    href = f'<a href="data:file/txt;base64,{b64}" download="{filename}">Download Summary</a>'
    st.markdown(href, unsafe_allow_html=True)

# Highlight differences between original and summarized text in blue
def highlight_differences(original, summarized):
    diff = difflib.ndiff(original.split(), summarized.split())
    diff_html = ''.join([f"<span style='color: lightblue;'>{word[2:]}</span> " if word.startswith('- ') 
                         else f"<span style='color: dodgerblue;'>{word[2:]}</span> " if word.startswith('+ ') 
                         else word[2:] + ' ' for word in diff])
    st.markdown(f"<div style='font-family: monospace;'>{diff_html}</div>", unsafe_allow_html=True)

# Load NLP tools for sentiment analysis
@st.cache_resource
def load_nlp_tools():
    # Using the nlptown/bert-base-multilingual-uncased-sentiment model for sentiment analysis
    sentiment_analyzer = pipeline('sentiment-analysis', model="nlptown/bert-base-multilingual-uncased-sentiment")
    return sentiment_analyzer

# Load custom CSS for styling
def local_css(file_name):
    with open(file_name) as f:
        st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)

local_css("style.css")

# Initialize session state variables if they don't exist yet
if "summary_text" not in st.session_state:
    st.session_state.summary_text = ""  # Initialize with an empty string

if "keywords" not in st.session_state:
    st.session_state.keywords = []  # Initialize with an empty list

# Load SpaCy model for keyword extraction
@st.cache_resource
def load_spacy_model():
    return spacy.load("en_core_web_sm")

# Extract keywords using SpaCy
def extract_keywords(text, max_keywords=5):
    nlp = load_spacy_model()
    doc = nlp(text)
    keywords = [chunk.text for chunk in doc.noun_chunks if chunk.root.pos_ in ['NOUN', 'PROPN']]
    return keywords[:max_keywords]

# Title and Sidebar
st.title("Text & Document Summarizer")
st.sidebar.image("logo_processed.png", use_column_width=True)

# Sidebar navigation and page choices
choice = st.sidebar.selectbox("Select your choice", ["Summarize Text", "Summarize Document"])

# Adding Sentiment Analysis Option in Sidebar
st.sidebar.subheader("Sentiment Analysis")
sentiment_analysis_enabled = st.sidebar.checkbox("Enable Sentiment Analysis")

# Adding Extract Keywords Option in Sidebar
st.sidebar.subheader("Extract Keywords")
extract_keywords_enabled = st.sidebar.checkbox("Enable Keyword Extraction")

# Adding Highlight Differences Option in Sidebar
st.sidebar.subheader("Text Comparison")
highlight_differences_enabled = st.sidebar.checkbox("Highlight Differences")

# Summarize Text page
if choice == "Summarize Text":
    st.subheader("Summarizing Text")

    # Input text area for summarizing
    input_text = st.text_area("Enter your text here")

    if input_text:
        input_word_count = len(input_text.split())
        st.write(f"Input Text Word Count: {input_word_count}")  # Display word count for input text

        # Perform keyword extraction if enabled
        if extract_keywords_enabled:
            keywords = extract_keywords(input_text)
            st.session_state.keywords = keywords  # Store extracted keywords in session state
            st.write("Extracted Keywords:", ', '.join(keywords))

        # Add slider to select the length of the summary
        summary_length = st.slider("Select summary length (max words)", min_value=50, max_value=500, step=1, value=100)

        # Perform normal summarization
        if st.button("Summarize Text"):
            start_time = time.time()  # Start timer before summarization
            summary = text_summary(input_text, maxlength=summary_length)
            end_time = time.time()  # End timer after summarization
            st.session_state.summary_text = summary  # Store summary in session state
            summary_word_count = len(summary.split())
            st.success(summary)
            st.write(f"Summary Word Count: {summary_word_count}")  # Display word count for summary

            time_taken = end_time - start_time  # Calculate time taken for summarization
            st.write(f"Time Taken to Summarize: {time_taken:.2f} seconds")  # Display time taken

            # Perform sentiment analysis if enabled
            if sentiment_analysis_enabled:
                sentiment_analyzer = load_nlp_tools()
                sentiment = sentiment_analyzer(st.session_state.summary_text)  # Analyze the summarized text
                sentiment_label = sentiment[0]['label']
                sentiment_score = sentiment[0]['score'] * 10  # Scale to be out of 10
                st.write(f"Sentiment Label: {sentiment_label}")
                st.write(f"Sentiment Score (out of 10): {sentiment_score:.2f}")

        # Highlight differences between original and summarized text
        if st.session_state.summary_text and highlight_differences_enabled:
            st.subheader("Highlight Differences")
            with st.expander("View Highlighted Differences", expanded=True):
                highlight_differences(input_text, st.session_state.summary_text)

        # Download Summary
        if st.button("Download Summary") and st.session_state.summary_text:
            download_summary(st.session_state.summary_text)

# Summarize Document page
elif choice == "Summarize Document":
    st.subheader("Summarizing Document")
    input_file = st.file_uploader("Upload your document here", type=['pdf', 'docx', 'jpg', 'png', 'pptx'])

    if input_file:
        file_type = input_file.name.split('.')[-1]
        extracted_text = ""

        if file_type == 'pdf':
            extracted_text = extract_text_from_pdf(input_file)
        elif file_type == 'docx':
            extracted_text = extract_text_from_word(input_file)
        elif file_type == 'pptx':
            extracted_text = extract_text_from_pptx(input_file)
        elif file_type in ['jpg', 'png']:
            extracted_text = extract_text_from_image(input_file)

        if extracted_text:
            with st.expander("Extracted Text", expanded=True):
                st.text_area("Extracted Text", value=extracted_text, height=200)

            extracted_word_count = len(extracted_text.split())
            st.write(f"Extracted Text Word Count: {extracted_word_count}")  # Word count for extracted text

            # Perform keyword extraction if enabled
            if extract_keywords_enabled:
                keywords = extract_keywords(extracted_text)
                st.session_state.keywords = keywords  # Store extracted keywords in session state
                st.write("Extracted Keywords:", ', '.join(keywords))

            # Add slider to select the length of the summary
            summary_length = st.slider("Select summary length (max words)", min_value=50, max_value=500, step=1, value=100)

            # Perform normal summarization for documents
            if st.button("Summarize Document"):
                start_time = time.time()  # Start timer before summarization
                summary = text_summary(extracted_text, maxlength=summary_length)
                end_time = time.time()  # End timer after summarization
                st.session_state.summary_text = summary  # Store summary in session state
                summary_word_count = len(summary.split())
                st.success(summary)
                st.write(f"Summary Word Count: {summary_word_count}")  # Display word count for summary

                time_taken = end_time - start_time  # Calculate time taken for summarization
                st.write(f"Time Taken to Summarize: {time_taken:.2f} seconds")  # Display time taken

                # Perform sentiment analysis if enabled
                if sentiment_analysis_enabled:
                    sentiment_analyzer = load_nlp_tools()
                    sentiment = sentiment_analyzer(st.session_state.summary_text)  # Analyze the summarized text
                    sentiment_label = sentiment[0]['label']
                    sentiment_score = sentiment[0]['score'] * 10  # Scale to be out of 10
                    st.write(f"Sentiment Label: {sentiment_label}")
                    st.write(f"Sentiment Score (out of 10): {sentiment_score:.2f}")

        # Highlight differences between original and summarized text
        if st.session_state.summary_text and highlight_differences_enabled:
            st.subheader("Highlight Differences")
            with st.expander("View Highlighted Differences", expanded=True):
                highlight_differences(extracted_text, st.session_state.summary_text)

        # Download Summary
        if st.button("Download Summary") and st.session_state.summary_text:
            download_summary(st.session_state.summary_text)